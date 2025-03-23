import fitz  # PyMuPDF is imported as fitz
import re
import docx
from pptx import Presentation
import openpyxl
import msoffcrypto

import fitz  # PyMuPDF
import re

def mask_pii_in_text(text):
    if text is None:
        return ""

    # Mask email addresses (show first 2 characters)
    text = re.sub(r'\b([A-Za-z0-9._%+-]{2})[A-Za-z0-9._%+-]*@([A-Za-z0-9.-]*?)([A-Za-z0-9]{2})(\.[A-Z|a-z]{2,7})\b', r'\1***@***\3\4', text)

    # Mask phone numbers (global format, show only last 4 digits)
    text = re.sub(r'\b(\+?\d{1,4}[-.\s]??\(?\d{1,4}\)?[-.\s]??\d{1,4}[-.\s]??\d{1,4}[-.\s]??)(\d{4})\b', r'XXX-XXX-\2', text)

    # Mask Aadhaar numbers (12-digit format, no leading 0 or 1, with spaces)
    text = re.sub(r'\b([2-9]{1}[0-9]{1})[0-9]{8}([0-9]{2})\b', r'\1XX-XXXX-XX\2', text)

    # Mask PAN card numbers (show first 3 letters, mask next 2 letters, show last 4 digits and last letter)
    text = re.sub(r'\b([A-Z]{2})[A-Z]{3}([0-9]{3})([0-9]{1}[A-Z]{1})\b', r'\1XXX\2X\3', text)

    # Mask Driving License numbers (show first 4 chars and last 4 digits)
    text = re.sub(r'\b([A-Z]{2}[0-9]{2})[0-9A-Z]{6,8}([0-9]{4})\b', r'\1-XXXXXX-\2', text)

    return text

def modify_and_encrypt_pdf(input_stream, output_stream, password):
    try:
        # Create a new PDF document from the stream
        doc = fitz.open(stream=input_stream.read(), filetype="pdf")
        
        # Process each page
        for page in doc:
            # Process text blocks instead of whole page text
            blocks = page.get_text("blocks")
            
            for block in blocks:
                block_text = block[4]
                masked_text = mask_pii_in_text(block_text)
                
                if masked_text != block_text:  # Only modify if there's a change
                    # Redact the original block (visually hide it)
                    page.add_redact_annot(block[:4], fill=(1, 1, 1))  # White fill for redaction
                    
            # Apply all redactions first
            page.apply_redactions()
            
            # Then insert masked text for each block that needed masking
            for block in blocks:
                block_text = block[4]
                masked_text = mask_pii_in_text(block_text)
                
                if masked_text != block_text:
                    # Insert the masked text at the same location
                    page.insert_text((block[0], block[1]), masked_text, fontsize=12)  # Adjust fontsize as needed
        
        # Save with encryption
        doc.save(
            output_stream,
            encryption=fitz.PDF_ENCRYPT_AES_256,  # Use AES 256-bit encryption
            user_pw=password,  # User password
            owner_pw=password,  # Owner password
            permissions=int(
                fitz.PDF_PERM_PRINT |  # Allow printing
                fitz.PDF_PERM_COPY    # Allow copying text
            )
        )
        doc.close()
        return True
    except Exception as e:
        print(f"Error in modify_and_encrypt_pdf: {str(e)}")
        raise

def modify_docx(input_stream, output_stream):
    doc = docx.Document(input_stream)

    # Iterate through each paragraph and replace PII
    for paragraph in doc.paragraphs:
        original_text = paragraph.text
        masked_text = mask_pii_in_text(original_text)
        if original_text != masked_text:
            paragraph.text = masked_text

    # Iterate through each table and replace PII
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                original_text = cell.text
                masked_text = mask_pii_in_text(original_text)
                if original_text != masked_text:
                    cell.text = masked_text

    # Save the modified document
    doc.save(output_stream)

def mask_pptx_file(input_stream, output_stream):
    prs = Presentation(input_stream)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        masked_text = mask_pii_in_text(original_text)
                        if original_text != masked_text:
                            run.text = masked_text

    prs.save(output_stream)

def mask_excel_file(input_stream, output_stream):
    wb = openpyxl.load_workbook(input_stream)

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    original_text = cell.value
                    masked_text = mask_pii_in_text(original_text)
                    if original_text != masked_text:
                        cell.value = masked_text

    wb.save(output_stream)

def encrypt_excel_file(input_stream):
    office_file = msoffcrypto.OfficeFile(input_stream)
    office_file.load_key(password="securepassword")
    office_file.encrypt("securepassword", input_stream)